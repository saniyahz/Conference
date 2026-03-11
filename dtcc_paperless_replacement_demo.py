from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ──
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
MEDIUM_BLUE = RGBColor(0x2E, 0x5C, 0x8A)
LIGHT_BLUE = RGBColor(0x3A, 0x7C, 0xBD)
ACCENT_ORANGE = RGBColor(0xE8, 0x6C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF4, 0xF7)
DARK_GRAY = RGBColor(0x4A, 0x4A, 0x4A)
RED_PAIN = RGBColor(0xC0, 0x39, 0x2B)
GREEN_NEW = RGBColor(0x27, 0xAE, 0x60)
SUBTLE_GRAY = RGBColor(0xE0, 0xE0, 0xE0)


def add_solid_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide_content(tf, items, font_size=18, color=DARK_GRAY, spacing=Pt(12)):
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 1 — Title Slide
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_solid_bg(slide, DARK_BLUE)

# Accent bar at top
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_ORANGE)

# Title
add_textbox(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.2),
            "DTCC Paperless Replacement", font_size=44, bold=True, color=WHITE)

# Subtitle
add_textbox(slide, Inches(1), Inches(3.4), Inches(11), Inches(0.8),
            "Automating the Ceding Carrier Replacement Process", font_size=24, color=LIGHT_BLUE)

# Divider line
div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(4.4), Inches(3), Inches(0.04))
div.fill.solid()
div.fill.fore_color.rgb = ACCENT_ORANGE
div.line.fill.background()

# Meta info
add_textbox(slide, Inches(1), Inches(4.7), Inches(11), Inches(0.5),
            "Feature Demo  |  March 2026", font_size=16, color=SUBTLE_GRAY)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 2 — Current State (Pain Points)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)

# Top accent bar
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_ORANGE)

# Section label
add_textbox(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.4),
            "CURRENT STATE", font_size=12, bold=True, color=ACCENT_ORANGE)

# Slide title
add_textbox(slide, Inches(0.8), Inches(0.75), Inches(11), Inches(0.7),
            "The Replacement Process Today — Fully Manual", font_size=32, bold=True, color=DARK_BLUE)

# Left column - process steps
left_box = add_shape_bg(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8), LIGHT_GRAY)
left_box.fill.solid()
left_box.fill.fore_color.rgb = LIGHT_GRAY

add_textbox(slide, Inches(1.1), Inches(1.95), Inches(5), Inches(0.4),
            "Manual Workflow Steps", font_size=16, bold=True, color=MEDIUM_BLUE)

steps = [
    "1.  Replacement comes into a case",
    "2.  Analyst manually identifies the ceding carrier",
    "3.  Analyst sends email, fax, or phone call\n     to the ceding carrier requesting documents",
    "4.  Wait for ceding carrier to respond",
    "5.  Manually track follow-ups & status",
    "6.  Manually enter received data into the system",
]

txBox = slide.shapes.add_textbox(Inches(1.1), Inches(2.5), Inches(5), Inches(3.8))
tf = txBox.text_frame
tf.word_wrap = True
for i, step in enumerate(steps):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = step
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_GRAY
    p.font.name = "Calibri"
    p.space_after = Pt(10)

# Right column - pain points
right_box = add_shape_bg(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), RGBColor(0xFD, 0xED, 0xED))

add_textbox(slide, Inches(7.1), Inches(1.95), Inches(5), Inches(0.4),
            "Pain Points", font_size=16, bold=True, color=RED_PAIN)

pains = [
    "\u2716  Time-consuming manual outreach (email / fax / calls)",
    "\u2716  No standardized process — varies by analyst",
    "\u2716  Difficult to track status of pending requests",
    "\u2716  Delays in receiving replacement paperwork",
    "\u2716  Risk of human error in data entry",
    "\u2716  No audit trail for communications",
]

txBox2 = slide.shapes.add_textbox(Inches(7.1), Inches(2.5), Inches(5.1), Inches(3.8))
tf2 = txBox2.text_frame
tf2.word_wrap = True
for i, pain in enumerate(pains):
    if i == 0:
        p = tf2.paragraphs[0]
    else:
        p = tf2.add_paragraph()
    p.text = pain
    p.font.size = Pt(16)
    p.font.color.rgb = RED_PAIN
    p.font.name = "Calibri"
    p.space_after = Pt(10)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 3 — New Feature: DTCC Paperless Replacement
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_ORANGE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.4),
            "NEW FEATURE", font_size=12, bold=True, color=GREEN_NEW)

add_textbox(slide, Inches(0.8), Inches(0.75), Inches(11), Inches(0.7),
            "DTCC Paperless Replacement — What's New", font_size=32, bold=True, color=DARK_BLUE)

# Feature cards
features = [
    {
        "title": "Auto-Select Owner on TRX Creation",
        "desc": "System will automatically select the owner while creating a transaction — no more manual lookup or selection errors.",
        "color": GREEN_NEW,
    },
    {
        "title": "Decimal Precision Fix",
        "desc": "Fixed the decimal issue — values now correctly support the format xxx.x0 (e.g., 125.50 instead of truncating to 125.5).",
        "color": LIGHT_BLUE,
    },
    {
        "title": "Paperless Communication via DTCC",
        "desc": "Replacement requests are sent electronically through DTCC — eliminating manual emails, faxes, and phone calls to ceding carriers.",
        "color": MEDIUM_BLUE,
    },
    {
        "title": "End-to-End Tracking",
        "desc": "Full visibility into replacement request status — sent, received, pending — with an auditable trail of all communications.",
        "color": ACCENT_ORANGE,
    },
]

card_width = Inches(5.5)
card_height = Inches(2.1)
x_positions = [Inches(0.8), Inches(7)]
y_positions = [Inches(1.8), Inches(4.2)]

for i, feat in enumerate(features):
    x = x_positions[i % 2]
    y = y_positions[i // 2]

    # Card background
    card = add_shape_bg(slide, x, y, card_width, card_height, LIGHT_GRAY)

    # Color accent bar on left of card
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), card_height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = feat["color"]
    accent.line.fill.background()

    # Feature title
    add_textbox(slide, x + Inches(0.3), y + Inches(0.2), card_width - Inches(0.5), Inches(0.5),
                feat["title"], font_size=20, bold=True, color=DARK_BLUE)

    # Feature description
    add_textbox(slide, x + Inches(0.3), y + Inches(0.75), card_width - Inches(0.5), Inches(1.2),
                feat["desc"], font_size=15, color=DARK_GRAY)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 4 — Before vs After
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)

add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_ORANGE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.4),
            "COMPARISON", font_size=12, bold=True, color=ACCENT_ORANGE)

add_textbox(slide, Inches(0.8), Inches(0.75), Inches(11), Inches(0.7),
            "Before vs. After", font_size=32, bold=True, color=DARK_BLUE)

# Before column
before_bg = add_shape_bg(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5), RGBColor(0xFD, 0xED, 0xED))

add_textbox(slide, Inches(1.1), Inches(1.95), Inches(5), Inches(0.5),
            "BEFORE  (Manual)", font_size=20, bold=True, color=RED_PAIN)

befores = [
    "\u2716  Manually identify ceding carrier",
    "\u2716  Email / fax / call to request docs",
    "\u2716  Manually select owner on each TRX",
    "\u2716  Decimal values truncated (xxx.x)",
    "\u2716  No visibility into request status",
    "\u2716  No audit trail",
    "\u2716  Slow turnaround, frequent follow-ups",
]

txBox = slide.shapes.add_textbox(Inches(1.1), Inches(2.6), Inches(5), Inches(3.8))
tf = txBox.text_frame
tf.word_wrap = True
for i, item in enumerate(befores):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = item
    p.font.size = Pt(16)
    p.font.color.rgb = RED_PAIN
    p.font.name = "Calibri"
    p.space_after = Pt(14)

# After column
after_bg = add_shape_bg(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5), RGBColor(0xE8, 0xF8, 0xF0))

add_textbox(slide, Inches(7.1), Inches(1.95), Inches(5), Inches(0.5),
            "AFTER  (DTCC Paperless)", font_size=20, bold=True, color=GREEN_NEW)

afters = [
    "\u2714  Automated electronic request via DTCC",
    "\u2714  No more emails, faxes, or phone calls",
    "\u2714  System auto-selects owner on TRX creation",
    "\u2714  Decimal precision fixed (xxx.x0)",
    "\u2714  Real-time status tracking",
    "\u2714  Full audit trail of all communications",
    "\u2714  Faster turnaround, fewer manual steps",
]

txBox2 = slide.shapes.add_textbox(Inches(7.1), Inches(2.6), Inches(5.1), Inches(3.8))
tf2 = txBox2.text_frame
tf2.word_wrap = True
for i, item in enumerate(afters):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    p.text = item
    p.font.size = Pt(16)
    p.font.color.rgb = GREEN_NEW
    p.font.name = "Calibri"
    p.space_after = Pt(14)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 5 — Demo / Thank You
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, DARK_BLUE)

add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT_ORANGE)

add_textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
            "Live Demo", font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(3.3), Inches(11), Inches(0.8),
            "DTCC Paperless Replacement in Action", font_size=24, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.3), Inches(2.3), Inches(0.04))
div.fill.solid()
div.fill.fore_color.rgb = ACCENT_ORANGE
div.line.fill.background()

add_textbox(slide, Inches(1), Inches(4.7), Inches(11), Inches(0.5),
            "Questions & Discussion", font_size=18, color=SUBTLE_GRAY, alignment=PP_ALIGN.CENTER)


# ── Save ──
output_path = "/home/user/Conference/DTCC_Paperless_Replacement_Demo.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
